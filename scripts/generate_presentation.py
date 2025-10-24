import json
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt, Inches
import pptx.api
import zipfile
import xml.etree.ElementTree as ET
import os
import requests
import tempfile
import shutil
import urllib.parse
import time

# Pārrakstām PPTX atpazīšanas loģiku, lai apstrādātu arī dažus citus content_type
def new_is_pptx_package(presentation_part):
    return presentation_part.content_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
        "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
    )

pptx.api._is_pptx_package = new_is_pptx_package

# Komandrindas parametri: template_path output_path data_path [unsplash_access_key]
if len(sys.argv) not in (4, 5):
    print("Usage: python generate_presentation.py <template_path> <output_path> <data_path> [unsplash_access_key]")
    sys.exit(1)

template_path = sys.argv[1]
output_path = sys.argv[2]
data_path = sys.argv[3]
unsplash_key = None
if len(sys.argv) == 5:
    unsplash_key = sys.argv[4]
    if sys.argv[4] else None

# Ja atslēga nav norādīta komandrindā, mēģinām to ņemt no vides mainīgā
if not unsplash_key:
    unsplash_key = os.environ.get('UNSPLASH_ACCESS_KEY')

# Iegūstam ievaddatus (JSON), kuros jābūt slides sarakstam vai 'topic'
with open(data_path, 'r') as f:
    data = json.load(f)

slides_data = data.get('slides', [])


# Helper funkcija: lejupielādē attēlu pēc meklējuma virsraksta
# Atgriež ceļu uz pagaidu failu vai None, ja neizdodas
def download_image_for_query(query, size=(1600, 900)):
    """Return local file path to downloaded image or None on failure."""
    safe_query = query.strip()
    if not safe_query:
        return None

    try:
        if unsplash_key:
            # Ja ir OAuth atslēga, izmanto oficiālo API meklēšanu (precīzāka)
            search_url = 'https://api.unsplash.com/search/photos'
            params = {
                'query': safe_query,
                'per_page': 1,
                'orientation': 'landscape',
            }
            headers = {
                'Accept-Version': 'v1',
                'Authorization': f'Client-ID {unsplash_key}'
            }
            r = requests.get(search_url, params=params, headers=headers, timeout=20)
            r.raise_for_status()
            body = r.json()
            results = body.get('results') or []
            if results:
                img_url = results[0].get('urls', {}).get('full') or results[0].get('urls', {}).get('regular')
            else:
                img_url = None
        else:
            # Ja nav atslēgas, izmanto public 'source.unsplash.com' ar query paramiem
            q = urllib.parse.quote_plus(safe_query)
            img_url = f'https://source.unsplash.com/{size[0]}x{size[1]}/?{q}'

        if not img_url:
            return None

        # Lejupielāde ar timeout un stream, saglabājam pagaidu JPG failā
        r = requests.get(img_url, stream=True, timeout=30)
        r.raise_for_status()

        suffix = '.jpg'
        tmp_fd, tmp_path = tempfile.mkstemp(suffix=suffix)
        os.close(tmp_fd)
        with open(tmp_path, 'wb') as out_f:
            shutil.copyfileobj(r.raw, out_f)

        return tmp_path
    except Exception as e:
        # Brīdinājums izdrukā uz stdout/stderr — noder, ja skriptu izsauc PHP
        print(f"[WARN] Image download failed for query '{query}': {e}")
        return None


# Atveram PowerPoint šablonu
prs = Presentation(template_path)

# Noņemam visus sākotnējos slaidus no šablona, lai strādātu ar tukšu prezentāciju
slide_ids = [sld.rId for sld in prs.slides._sldIdLst]
for rId in slide_ids:
    prs.part.drop_rel(rId)
prs.slides._sldIdLst.clear()

slide_layouts = prs.slide_layouts
title_layout = slide_layouts[0]
content_layout = slide_layouts[1]

tmp_images = []  # saraksts uz laika failiem, ko pēc tam dzēsīsim

# Izveidojam slaidus pēc JSON datiem
for i, slide_data in enumerate(slides_data):
    # Izmantosim title layout pirmajam slaidam, pārējiem content layout
    layout = title_layout if i == 0 else content_layout

    slide = prs.slides.add_slide(layout)

    # Attēls nav paredzēts pirmajam (title) slaidam
    place_image = (i != 0)

    # Noteiksim meklēšanas vaicājumu attēlam: prioritāte virsraksts + pirmais bullet
    query = ''
    if isinstance(slide_data, dict):
        title_text = slide_data.get('title', '') or ''
        bullets = slide_data.get('bullets', []) or []
    else:
        title_text = str(slide_data)
        bullets = []

    if not query:
        first_bullet = bullets[0] if bullets else ''
        query = (title_text + ' ' + first_bullet).strip()

    if not query:
        # Ja nav individuālā vaicājuma, izmanto prezentācijas 'topic' kā fallback
        query = data.get('topic', '')

    img_path = None
    if query and place_image:
        # Mēģinām lejupielādēt attēlu līdz 3 reizēm ar nelielu pauzi starp mēģinājumiem
        for attempt in range(1, 4):
            img_path = download_image_for_query(query)
            if img_path:
                tmp_images.append(img_path)
                break
            time.sleep(0.5 * attempt)

    if img_path:
        try:
            # Aprēķinām attēla izvietojumu: labajā puse, ar nelielu marginu
            margin = Inches(0.3)
            max_width = int(prs.slide_width * 0.4 - margin)
            left = int(prs.slide_width - max_width - margin)

            pic = slide.shapes.add_picture(img_path, left, int(margin), width=max_width)

            # Vertikāli centrujam attēlu slaidā
            pic.top = int((prs.slide_height - pic.height) / 2)

            # Pārvietojam visus tekstu elementus uz priekšplānu, lai teksts būtu redzams virs attēla
            text_elements = []
            for shp in list(slide.shapes):
                try:
                    if shp.has_text_frame and (not shp.is_placeholder or shp.placeholder_format.type != 1):
                        text_elements.append(shp._element)
                except Exception:
                    pass

            spTree = slide.shapes._spTree
            for el in text_elements:
                try:
                    spTree.remove(el)
                    spTree.append(el)
                except Exception:
                    pass

        except Exception as e:
            # Brīdinājums, ja attēla ievietošana neizdodas — skripts turpina darbu bez attēla
            print(f"[WARN] Failed to place image on slide {i}: {e}")

    # Uzstādām virsrakstu, ja tas ir pieejams un definēts JSON
    title_shape = None
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None

    if title_shape and isinstance(slide_data, dict):
        if slide_data.get('title'):
            title_shape.text = slide_data['title']

    # Mēģinām atrast body placeholder (BODY vai SUBTITLE uz pirmajiem slaidiem)
    body_shape = None
    try:
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            if ph_type == 7:  # BODY
                body_shape = shape
                break
            elif i == 0 and ph_type == 17:  # SUBTITLE for title slide
                body_shape = shape
                break
    except Exception:
        body_shape = None

    # Fallback: ja nav placeholder, izvēlamies lielāko tekstu laukumu
    if body_shape is None:
        candidates = [sh for sh in slide.shapes if sh.has_text_frame and sh != title_shape]
        if candidates:
            body_shape = max(candidates, key=lambda sh: sh.height)

    # Aizpildām bullet punktus, ja tie definēti JSON
    if body_shape and isinstance(slide_data, dict) and 'bullets' in slide_data:
        bullets = slide_data['bullets']
        tf = body_shape.text_frame
        if bullets:
            # Uzstādām pirmo rindu uz pirmo bullet, pārējie kā papildu rindas
            tf.text = bullets[0]
            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
            # Fiksējam fonta izmēru visiem runiem, lai izskatītos konsekventi
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(18)
        else:
            tf.clear()

# Saglabājam izveidoto prezentāciju uz diska
prs.save(output_path)

# Dzēšam visus pagaidu attēlu failus
for p in tmp_images:
    try:
        os.remove(p)
    except Exception:
        pass

# Labojums [Content_Types].xml — nodrošinām pareizu ContentType ierakstu prezentācijas daļai
# Dažos vides apstākļos tas var būt nepieciešams, lai Office atpazītu failu kā .pptx
temp_path = output_path + '.temp'
with zipfile.ZipFile(output_path, 'r') as zin:
    with zipfile.ZipFile(temp_path, 'w') as zout:
        for item in zin.infolist():
            if item.filename == '[Content_Types].xml':
                xml = zin.read(item.filename)
                tree = ET.fromstring(xml)
                for elem in tree.iter():
                    if 'PartName' in elem.attrib and elem.attrib['PartName'] == '/ppt/presentation.xml':
                        elem.attrib['ContentType'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
                new_xml = ET.tostring(tree, encoding='utf-8', method='xml')
                zout.writestr(item, new_xml)
            else:
                zout.writestr(item, zin.read(item.filename))
os.remove(output_path)
os.rename(temp_path, output_path)

print(f"Presentation saved to {output_path}")
